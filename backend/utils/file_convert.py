"""
文件转换核心模块 - 在线文件转换工具
支持 Word、PDF、图片等常见办公文件的格式转换
"""
import os
import uuid


def word_to_pdf(input_path, output_path):
    """
    Word转PDF
    使用LibreOffice进行转换（需要系统安装LibreOffice）
    
    Args:
        input_path: 输入Word文件路径
        output_path: 输出PDF文件路径
    
    Returns:
        (success, path, message) 元组
    """
    try:
        import subprocess
        import shutil
        
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', os.path.dirname(output_path),
            input_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        # LibreOffice会在同目录生成同名.pdf文件
        expected = os.path.join(
            os.path.dirname(output_path), 
            os.path.splitext(os.path.basename(input_path))[0] + '.pdf'
        )
        
        if os.path.exists(expected):
            if expected != output_path:
                shutil.move(expected, output_path)
            return True, output_path, "转换成功"
        else:
            return False, None, f"转换失败: {result.stderr or '文件未生成'}"
            
    except subprocess.TimeoutExpired:
        return False, None, "转换超时"
    except Exception as e:
        return False, None, f"Word转PDF失败: {str(e)}"


def pdf_to_word(input_path, output_path):
    """
    PDF转Word
    使用pdf2docx库进行转换
    
    Args:
        input_path: 输入PDF文件路径
        output_path: 输出Word文件路径
    
    Returns:
        (success, path, message) 元组
    """
    try:
        from pdf2docx import Converter
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        return True, output_path, "转换成功"
    except Exception as e:
        return False, None, f"PDF转Word失败: {str(e)}"


def pdf_to_pptx(input_path, output_path):
    """
    PDF转PPT
    - 文字类PDF直接提取文字
    - 图片类PDF（扫描件）使用OCR识别
    
    Args:
        input_path: 输入PDF文件路径
        output_path: 输出PPT文件路径
    
    Returns:
        (success, path, message) 元组
    """
    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import pytesseract
        from PIL import Image
        import io
        
        doc = fitz.open(input_path)
        prs = Presentation()
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text().strip()
            
            # 如果没有提取到文字，尝试OCR识别
            if not text:
                try:
                    # 将PDF页面渲染为图片
                    mat = fitz.Matrix(2, 2)  # 2x分辨率
                    pix = page.get_pixmap(matrix=mat)
                    img_data = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_data))
                    
                    # OCR识别（支持中文和英文）
                    text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                    img.close()
                except Exception as ocr_err:
                    print(f"OCR识别失败: {ocr_err}")
                    text = ""
            
            # 创建幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(6)
            
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            p = text_frame.paragraphs[0]
            if text:
                p.text = f"--- 第 {page_num + 1} 页 ---\n\n{text[:5000]}"
            else:
                p.text = f"--- 第 {page_num + 1} 页 ---\n\n[此页无法识别内容]"
            p.font.size = Pt(11)
        
        doc.close()
        prs.save(output_path)
        return True, output_path, "转换成功"
    except Exception as e:
        return False, None, f"PDF转PPT失败: {str(e)}"


def images_to_pdf(input_paths, output_path):
    """
    多张图片合并为一个PDF
    
    Args:
        input_paths: 输入图片文件路径列表
        output_path: 输出PDF文件路径
    
    Returns:
        (success, path, message) 元组
    """
    try:
        from PIL import Image
        
        images = []
        for path in input_paths:
            img = Image.open(path)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            images.append(img)
        
        if len(images) == 1:
            images[0].save(output_path, "PDF")
        else:
            images[0].save(output_path, "PDF", save_all=True, append_images=images[1:])
        
        for img in images:
            img.close()
        return True, output_path, "转换成功"
    except Exception as e:
        return False, None, f"图片转PDF失败: {str(e)}"


def pdf_to_images(input_path, output_dir):
    """
    PDF转图片（返回多张图片路径）
    
    Args:
        input_path: 输入PDF文件路径
        output_dir: 输出目录
    
    Returns:
        (success, paths_list, message) 元组
    """
    try:
        import fitz
        import os
        
        doc = fitz.open(input_path)
        output_paths = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            
            output_path = os.path.join(output_dir, f"page_{page_num + 1}.png")
            pix.save(output_path)
            output_paths.append(output_path)
        
        doc.close()
        return True, output_paths, "转换成功"
    except Exception as e:
        return False, None, f"PDF转图片失败: {str(e)}"


def merge_pdfs(input_paths, output_path):
    """
    合并多个PDF为一个
    
    Args:
        input_paths: 输入PDF文件路径列表
        output_path: 输出PDF文件路径
    
    Returns:
        (success, path, message) 元组
    """
    try:
        import fitz
        
        writer = fitz.open()
        for path in input_paths:
            reader = fitz.open(path)
            writer.insert_pdf(reader)
            reader.close()
        
        writer.save(output_path)
        writer.close()
        return True, output_path, "合并成功"
    except Exception as e:
        return False, None, f"合并PDF失败: {str(e)}"


def split_pdf(input_path, output_dir, page_ranges):
    """
    拆分PDF
    
    Args:
        input_path: 输入PDF文件路径
        output_dir: 输出目录
        page_ranges: 页码范围列表，如 [(1,3), (5,7)] 表示第1-3页和第5-7页
    
    Returns:
        (success, paths_list, message) 元组
    """
    try:
        import fitz
        import os
        
        doc = fitz.open(input_path)
        output_paths = []
        
        for i, (start, end) in enumerate(page_ranges):
            new_doc = fitz.open()
            new_doc.insert_pdf(doc, from_page=start-1, to_page=end-1)
            
            output_path = os.path.join(output_dir, f"split_{i+1}.pdf")
            new_doc.save(output_path)
            new_doc.close()
            output_paths.append(output_path)
        
        doc.close()
        return True, output_paths, "拆分成功"
    except Exception as e:
        return False, None, f"拆分PDF失败: {str(e)}"


def compress_pdf(input_path, output_path):
    """
    压缩PDF
    
    Args:
        input_path: 输入PDF文件路径
        output_path: 输出PDF文件路径
    
    Returns:
        (success, path, message) 元组
    """
    try:
        import fitz
        
        doc = fitz.open(input_path)
        
        for page in doc:
            page.clean_contents()
        
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        doc.close()
        return True, output_path, "压缩成功"
    except Exception as e:
        return False, None, f"压缩PDF失败: {str(e)}"


def convert_file(file_path, convert_type, output_dir):
    """
    单文件转换
    
    Args:
        file_path: 输入文件路径
        convert_type: 转换类型
        output_dir: 输出目录
    
    Returns:
        (success, path_or_list, message) 元组
    """
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    unique_id = uuid.uuid4().hex[:8]
    
    if convert_type == 'word2pdf':
        output_path = os.path.join(output_dir, f"{unique_id}_{base_name}.pdf")
        return word_to_pdf(file_path, output_path)
    
    elif convert_type == 'pdf2word':
        output_path = os.path.join(output_dir, f"{unique_id}_{base_name}.docx")
        return pdf_to_word(file_path, output_path)
    
    elif convert_type == 'pdf2pptx':
        output_path = os.path.join(output_dir, f"{unique_id}_{base_name}.pptx")
        return pdf_to_pptx(file_path, output_path)
    
    elif convert_type in ['jpg2pdf', 'png2pdf']:
        output_path = os.path.join(output_dir, f"{unique_id}_{base_name}.pdf")
        return images_to_pdf([file_path], output_path)
    
    elif convert_type == 'pdf2images':
        return pdf_to_images(file_path, output_dir)
    
    elif convert_type == 'compress':
        output_path = os.path.join(output_dir, f"{unique_id}_{base_name}_compressed.pdf")
        return compress_pdf(file_path, output_path)
    
    else:
        return False, None, f"不支持的转换类型: {convert_type}"


def convert_multiple_files(file_paths, convert_type, output_dir):
    """
    多文件转换（合并操作）
    
    Args:
        file_paths: 输入文件路径列表
        convert_type: 转换类型
        output_dir: 输出目录
    
    Returns:
        (success, path_or_list, message) 元组
    """
    unique_id = uuid.uuid4().hex[:8]
    
    if convert_type in ['jpg2pdf', 'png2pdf']:
        output_path = os.path.join(output_dir, f"{unique_id}_images.pdf")
        return images_to_pdf(file_paths, output_path)
    
    elif convert_type == 'merge':
        output_path = os.path.join(output_dir, f"{unique_id}_merged.pdf")
        return merge_pdfs(file_paths, output_path)
    
    else:
        return False, None, f"不支持的转换类型: {convert_type}"
